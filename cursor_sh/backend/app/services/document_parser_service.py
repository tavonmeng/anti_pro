"""客户资料解析服务。

支持可提取文本的 PDF / DOC / DOCX / PPTX / XLS / XLSX，不做 OCR。
"""

import os
import subprocess
from dataclasses import dataclass


class DocumentParseError(Exception):
    """文档解析失败。"""


@dataclass
class ParsedSection:
    label: str
    page: int | None
    text: str


MAX_SECTION_CHARS = 3500
MAX_TOTAL_CHARS = 60000
LEGACY_DOC_PARSE_TIMEOUT_SECONDS = 60
MAX_SPREADSHEET_COLUMNS = 100


def parse_document(file_path: str, filename: str = "") -> list[ParsedSection]:
    """解析文档文本并保留页码/幻灯片编号。"""
    ext = os.path.splitext(filename or file_path)[1].lower()
    if ext == ".pdf":
        return _parse_pdf(file_path)
    if ext == ".docx":
        return _parse_docx(file_path)
    if ext == ".doc":
        return _parse_doc(file_path)
    if ext == ".pptx":
        return _parse_pptx(file_path)
    if ext == ".xlsx":
        return _parse_xlsx(file_path)
    if ext == ".xls":
        return _parse_xls(file_path)
    raise DocumentParseError(f"不支持的文件类型: {ext}")


def build_llm_text(sections: list[ParsedSection]) -> str:
    """构建带来源标记的 LLM 输入文本。"""
    parts = []
    total = 0
    for section in sections:
        text = _clean_text(section.text)
        if not text:
            continue
        if len(text) > MAX_SECTION_CHARS:
            text = text[:MAX_SECTION_CHARS] + "\n...(本页内容已截断)"
        block = f"【来源：{section.label}】\n{text}"
        if total + len(block) > MAX_TOTAL_CHARS:
            break
        parts.append(block)
        total += len(block)
    return "\n\n".join(parts)


def build_llm_text_chunks(
    sections: list[ParsedSection],
    *,
    max_chunk_chars: int = 18000,
    max_total_chars: int = 120000,
) -> list[str]:
    """构建可分段送入 LLM 的带来源文本块。"""
    chunks: list[str] = []
    current_parts: list[str] = []
    current_len = 0
    total = 0
    max_chunk_chars = max(4000, int(max_chunk_chars or 18000))
    max_total_chars = max(max_chunk_chars, int(max_total_chars or 120000))

    for section in sections:
        text = _clean_text(section.text)
        if not text:
            continue
        if len(text) > MAX_SECTION_CHARS:
            text = text[:MAX_SECTION_CHARS] + "\n...(本页内容已截断)"
        block = f"【来源：{section.label}】\n{text}"
        if total + len(block) > max_total_chars:
            break
        if current_parts and current_len + len(block) > max_chunk_chars:
            chunks.append("\n\n".join(current_parts))
            current_parts = []
            current_len = 0
        current_parts.append(block)
        current_len += len(block)
        total += len(block)

    if current_parts:
        chunks.append("\n\n".join(current_parts))
    return chunks


def _parse_pdf(file_path: str) -> list[ParsedSection]:
    try:
        from pypdf import PdfReader
    except Exception as exc:
        raise DocumentParseError("缺少 PDF 解析依赖 pypdf，请安装 requirements.txt") from exc

    try:
        reader = PdfReader(file_path)
        sections = []
        for idx, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            sections.append(ParsedSection(label=f"第{idx}页", page=idx, text=text))
        return sections
    except Exception as exc:
        raise DocumentParseError(f"PDF 解析失败: {exc}") from exc


def _parse_docx(file_path: str) -> list[ParsedSection]:
    try:
        from docx import Document
    except Exception as exc:
        raise DocumentParseError("缺少 Word 解析依赖 python-docx，请安装 requirements.txt") from exc

    try:
        doc = Document(file_path)
        lines = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    lines.append(" | ".join(cells))
        return [ParsedSection(label="Word正文", page=None, text="\n".join(lines))]
    except Exception as exc:
        raise DocumentParseError(f"Word 解析失败: {exc}") from exc


def _parse_doc(file_path: str) -> list[ParsedSection]:
    """Extract text from a legacy Word 97-2003 binary document."""
    try:
        result = subprocess.run(
            ["antiword", "-m", "UTF-8.txt", file_path],
            check=False,
            capture_output=True,
            timeout=LEGACY_DOC_PARSE_TIMEOUT_SECONDS,
        )
    except FileNotFoundError as exc:
        raise DocumentParseError("缺少旧版 Word 解析依赖 antiword") from exc
    except subprocess.TimeoutExpired as exc:
        raise DocumentParseError("旧版 Word 解析超时") from exc
    except OSError as exc:
        raise DocumentParseError("旧版 Word 解析器启动失败") from exc

    if result.returncode != 0:
        error = result.stderr.decode("utf-8", errors="replace").strip()
        detail = f": {error[:240]}" if error else ""
        raise DocumentParseError(f"旧版 Word 解析失败{detail}")

    text = result.stdout.decode("utf-8", errors="replace")
    text = text.replace("\x00", "").strip()
    if not text:
        raise DocumentParseError("旧版 Word 没有可提取的文字")
    return [ParsedSection(label="Word正文", page=None, text=text)]


def _parse_pptx(file_path: str) -> list[ParsedSection]:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise DocumentParseError("缺少 PPT 解析依赖 python-pptx，请安装 requirements.txt") from exc

    try:
        prs = Presentation(file_path)
        sections = []
        for idx, slide in enumerate(prs.slides, start=1):
            lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    lines.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            lines.append(" | ".join(cells))
            sections.append(ParsedSection(label=f"第{idx}页幻灯片", page=idx, text="\n".join(lines)))
        return sections
    except Exception as exc:
        raise DocumentParseError(f"PPT 解析失败: {exc}") from exc


def _parse_xlsx(file_path: str) -> list[ParsedSection]:
    """Extract readable cell values from each worksheet in a modern Excel file."""
    try:
        from openpyxl import load_workbook
    except Exception as exc:
        raise DocumentParseError("缺少 Excel 解析依赖 openpyxl，请安装 requirements.txt") from exc

    try:
        workbook = load_workbook(file_path, read_only=True, data_only=True)
        sections = []
        for sheet in workbook.worksheets:
            rows = _spreadsheet_rows(sheet.iter_rows(values_only=True))
            sections.append(ParsedSection(label=f"工作表：{sheet.title}", page=None, text=rows))
        return sections
    except Exception as exc:
        raise DocumentParseError(f"Excel 解析失败: {exc}") from exc


def _parse_xls(file_path: str) -> list[ParsedSection]:
    """Extract readable cell values from a legacy Excel file."""
    try:
        import xlrd
    except Exception as exc:
        raise DocumentParseError("缺少旧版 Excel 解析依赖 xlrd，请安装 requirements.txt") from exc

    try:
        workbook = xlrd.open_workbook(file_path, on_demand=True)
        sections = []
        for sheet_name in workbook.sheet_names():
            sheet = workbook.sheet_by_name(sheet_name)
            rows = _spreadsheet_rows(
                (sheet.row_values(row_index) for row_index in range(sheet.nrows))
            )
            sections.append(ParsedSection(label=f"工作表：{sheet_name}", page=None, text=rows))
        return sections
    except Exception as exc:
        raise DocumentParseError(f"旧版 Excel 解析失败: {exc}") from exc


def _spreadsheet_rows(rows) -> str:
    """Serialize non-empty cells in a tabular, LLM-friendly form."""
    lines: list[str] = []
    for row in rows:
        values = [_spreadsheet_cell_text(cell) for cell in row[:MAX_SPREADSHEET_COLUMNS]]
        if any(values):
            lines.append(" | ".join(values))
    return "\n".join(lines)


def _spreadsheet_cell_text(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()


def _clean_text(text: str) -> str:
    lines = [line.strip() for line in (text or "").splitlines()]
    return "\n".join(line for line in lines if line)
