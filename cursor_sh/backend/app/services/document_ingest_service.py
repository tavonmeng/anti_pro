"""管理员设计方案文档 ingest 服务。

当前用于解析管理员上传到 AI 方案设计里的 PDF/PPTX/TXT 文件：
文件读取 -> 文本抽取 -> Qwen 结构化提取 -> 写回 orders.design_plan.files。
"""

from __future__ import annotations

import io
import json
import os
import re
from datetime import datetime, timezone
from typing import Any, Awaitable, Callable
from urllib.parse import unquote, urlparse

import httpx
from sqlalchemy import select
from sqlalchemy.orm.attributes import flag_modified

from app.config import settings
from app.database import async_session_maker
from app.models.order import Order


DOCUMENT_EXTENSIONS = {".pdf", ".pptx", ".txt", ".md"}
DOCUMENT_CHUNK_PAGES = 6
DOCUMENT_DIRECT_MAX_CHARS = 60000


def is_ingestable_file(file_info: dict[str, Any]) -> bool:
    """判断设计方案文件是否可进入文档 ingest。"""
    filename = str(file_info.get("filename") or file_info.get("name") or "")
    ext = os.path.splitext(filename)[1].lower()
    return ext in DOCUMENT_EXTENSIONS


async def ingest_design_plan_file(order_id: str, file_index: int, force: bool = False) -> None:
    """后台执行单个设计方案文件的 ingest。"""
    try:
        order, file_info = await _load_design_plan_file(order_id, file_index)
        if not order or not file_info:
            return

        if file_info.get("ingest_status") == "success" and not force:
            return

        await _update_file_ingest_state(
            order_id,
            file_index,
            {
                "ingest_status": "processing",
                "ingest_error": "",
                "ingest_started_at": _now_iso(),
                "ingest_model": settings.DOCUMENT_INGEST_MODEL_NAME,
            },
        )

        contents = await _load_file_bytes(file_info)
        filename = str(file_info.get("filename") or file_info.get("name") or "document")
        pages = _extract_document_pages(filename, contents)
        if not pages:
            raise ValueError("未能从文档中提取到可读文本，可能是扫描件或纯图片版 PDF/PPT")

        raw_text = _format_pages_for_llm(pages)
        result = await _extract_document_with_qwen(filename, pages)
        if not result:
            raise ValueError("Qwen 未返回有效的结构化解析结果")

        await _update_file_ingest_state(
            order_id,
            file_index,
            {
                "ingest_status": "success",
                "ingest_error": "",
                "ingested_at": _now_iso(),
                "ingest_model": settings.DOCUMENT_INGEST_MODEL_NAME,
                "ingest_page_count": len(pages),
                "ingest_text_chars": sum(len(p["text"]) for p in pages),
                "ingest_text_preview": raw_text[:3000],
                "ingest_result": result,
            },
        )

        try:
            from app.services.memory_service import sync_document_ingest

            await sync_document_ingest(order.user_id, {
                "order_id": order_id,
                "file_index": file_index,
                "filename": filename,
                **result,
            })
            await _update_file_ingest_state(
                order_id,
                file_index,
                {"memory_sync_status": "success", "memory_synced_at": _now_iso()},
            )
        except Exception as memory_exc:
            await _update_file_ingest_state(
                order_id,
                file_index,
                {
                    "memory_sync_status": "failed",
                    "memory_sync_error": str(memory_exc)[:500],
                },
            )
    except Exception as exc:
        await _update_file_ingest_state(
            order_id,
            file_index,
            {
                "ingest_status": "failed",
                "ingest_error": str(exc)[:500],
                "ingested_at": _now_iso(),
                "ingest_model": settings.DOCUMENT_INGEST_MODEL_NAME,
            },
        )


async def ingest_customer_document_bytes(
    user_id: str,
    document_id: str,
    filename: str,
    contents: bytes,
    source: str = "customer_profile_upload",
) -> None:
    """后台执行客户维度资料 ingest，结果直接同步到 UserMemory。"""
    from app.services.memory_service import sync_document_ingest, upsert_customer_document_ingest

    try:
        await upsert_customer_document_ingest(user_id, document_id, {
            "filename": filename,
            "source": source,
            "ingest_status": "processing",
            "ingest_error": "",
            "ingest_started_at": _now_iso(),
            "ingest_model": settings.DOCUMENT_INGEST_MODEL_NAME,
        })

        pages = _extract_document_pages(filename, contents)
        if not pages:
            raise ValueError("未能从文档中提取到可读文本，可能是扫描件或纯图片版 PDF/PPT")

        result = await _extract_document_with_qwen(filename, pages)
        if not result:
            raise ValueError("Qwen 未返回有效的结构化解析结果")

        await sync_document_ingest(user_id, {
            "document_id": document_id,
            "source": source,
            "filename": filename,
            **result,
        })
        await upsert_customer_document_ingest(user_id, document_id, {
            "filename": filename,
            "source": source,
            "ingest_status": "success",
            "ingest_error": "",
            "ingested_at": _now_iso(),
            "ingest_model": settings.DOCUMENT_INGEST_MODEL_NAME,
            "ingest_page_count": len(pages),
            "ingest_text_chars": sum(len(p["text"]) for p in pages),
            "ingest_result": result,
        })
    except Exception as exc:
        await upsert_customer_document_ingest(user_id, document_id, {
            "filename": filename,
            "source": source,
            "ingest_status": "failed",
            "ingest_error": str(exc)[:500],
            "ingested_at": _now_iso(),
            "ingest_model": settings.DOCUMENT_INGEST_MODEL_NAME,
        })


async def ingest_company_document_bytes(
    document_id: str,
    filename: str,
    contents: bytes,
    source: str = "company_profile_upload",
    raw_file: dict[str, Any] | None = None,
    progress_callback: Callable[[dict[str, Any]], Awaitable[None]] | None = None,
) -> dict[str, Any]:
    """管理员全局上传公司资料：自动识别公司名，归档到公司资料库。"""
    from app.services.company_profile_service import (
        link_profile_to_matching_users,
        sync_company_profile_ingest,
        update_company_library_document,
    )
    from app.services.company_library_storage import store_company_library_asset

    await update_company_library_document(document_id, {
        "status": "processing",
        "error": "",
    })
    pages = _extract_document_pages(filename, contents)
    if not pages:
        raise ValueError("未能从文档中提取到可读文本，可能是扫描件或纯图片版 PDF/PPT")

    text_chars = sum(len(p["text"]) for p in pages)
    raw_text = _format_pages_for_llm(pages)
    text_asset = store_company_library_asset(
        document_id=document_id,
        stage="extracted_text",
        filename=f"{_filename_stem(filename)}.md",
        data=raw_text.encode("utf-8"),
        content_type="text/markdown; charset=utf-8",
    )
    await update_company_library_document(document_id, {
        "status": "text_extracted",
        "page_count": str(len(pages)),
        "text_chars": str(text_chars),
        "extracted_text": text_asset,
        "text_preview": raw_text[:4000],
    })
    if progress_callback:
        await progress_callback({
            "stage": "text_extracted",
            "page_count": len(pages),
            "text_chars": text_chars,
            "extracted_text": text_asset,
        })

    result = await _extract_document_with_qwen(
        filename,
        pages,
        progress_callback=progress_callback,
    )
    if not result:
        raise ValueError("Qwen 未返回有效的结构化解析结果")

    company_info = result.get("company_info") or {}
    company_name = company_info.get("company_name") or ""
    if not company_name:
        raise ValueError("未能从资料中识别出客户公司名称")

    memory_asset = store_company_library_asset(
        document_id=document_id,
        stage="structured_memory",
        filename=f"{_filename_stem(filename)}.memory.json",
        data=json.dumps(result, ensure_ascii=False, indent=2).encode("utf-8"),
        content_type="application/json; charset=utf-8",
    )
    assets = {
        "raw_file": raw_file or {},
        "extracted_text": text_asset,
        "structured_memory": memory_asset,
    }
    profile = await sync_company_profile_ingest({
        "document_id": document_id,
        "source": source,
        "filename": filename,
        "assets": assets,
        "page_count": len(pages),
        "text_chars": text_chars,
        **result,
    })
    await update_company_library_document(document_id, {
        "status": "success",
        "error": "",
        "company_key": profile.company_key,
        "company_name": profile.company_name,
        "page_count": str(len(pages)),
        "text_chars": str(text_chars),
        "structured_memory": memory_asset,
    })
    matched_users = await link_profile_to_matching_users(company_name)
    return {
        "document_id": document_id,
        "company_key": profile.company_key,
        "company_name": profile.company_name,
        "matched_users": matched_users,
        "page_count": len(pages),
        "text_chars": text_chars,
        "assets": assets,
        "result": result,
    }


async def ingest_company_document_job(
    document_id: str,
    filename: str,
    contents: bytes,
    source: str = "company_profile_upload",
    raw_file: dict[str, Any] | None = None,
) -> None:
    """后台执行全局公司资料 ingest，并把进度写入任务表。"""
    from app.services.company_profile_service import (
        update_company_library_document,
        update_company_profile_ingest_job,
    )

    try:
        await update_company_profile_ingest_job(document_id, {
            "status": "processing",
            "error": "",
            "started_at": datetime.now(timezone.utc),
        })
        result = await ingest_company_document_bytes(
            document_id=document_id,
            filename=filename,
            contents=contents,
            source=source,
            raw_file=raw_file,
            progress_callback=lambda updates: update_company_profile_ingest_job(
                document_id,
                _job_progress_updates(updates),
            ),
        )
        await update_company_profile_ingest_job(document_id, {
            "status": "success",
            "error": "",
            "company_key": result.get("company_key", ""),
            "company_name": result.get("company_name", ""),
            "page_count": str(result.get("page_count") or ""),
            "text_chars": str(result.get("text_chars") or ""),
            "result": result,
            "finished_at": datetime.now(timezone.utc),
        })
    except Exception as exc:
        await update_company_library_document(document_id, {
            "status": "failed",
            "error": str(exc)[:1000],
        })
        await update_company_profile_ingest_job(document_id, {
            "status": "failed",
            "error": str(exc)[:1000],
            "finished_at": datetime.now(timezone.utc),
        })
        print(f"公司资料解析任务失败 {document_id} {filename}: {exc}")


async def _load_design_plan_file(order_id: str, file_index: int) -> tuple[Order | None, dict[str, Any] | None]:
    async with async_session_maker() as session:
        result = await session.execute(select(Order).where(Order.id == order_id))
        order = result.scalar_one_or_none()
        if not order:
            return None, None

        files = (order.design_plan or {}).get("files") or []
        if file_index < 0 or file_index >= len(files):
            return order, None
        return order, dict(files[file_index] or {})


async def _update_file_ingest_state(order_id: str, file_index: int, updates: dict[str, Any]) -> None:
    async with async_session_maker() as session:
        result = await session.execute(select(Order).where(Order.id == order_id))
        order = result.scalar_one_or_none()
        if not order:
            return

        plan = order.design_plan or {}
        files = plan.get("files") or []
        if file_index < 0 or file_index >= len(files):
            return

        file_info = dict(files[file_index] or {})
        file_info.update(updates)
        files[file_index] = file_info
        plan["files"] = files
        plan["updatedAt"] = _now_iso()
        order.design_plan = plan
        flag_modified(order, "design_plan")
        await session.commit()


async def _load_file_bytes(file_info: dict[str, Any]) -> bytes:
    url = str(file_info.get("url") or file_info.get("file_url") or "")
    local_path = _local_path_from_upload_url(url)
    if local_path:
        with open(local_path, "rb") as fh:
            return fh.read()

    if not url:
        raise ValueError("文件缺少 url，无法读取")

    async with httpx.AsyncClient(follow_redirects=True, timeout=60.0) as client:
        response = await client.get(url)
        response.raise_for_status()
        return response.content


def _local_path_from_upload_url(url: str) -> str | None:
    parsed = urlparse(url)
    path = unquote(parsed.path or url)
    if not path.startswith("/uploads/"):
        return None

    relative = path.removeprefix("/uploads/").lstrip("/")
    upload_root = os.path.abspath(settings.UPLOAD_DIR)
    file_path = os.path.abspath(os.path.join(upload_root, relative))
    if file_path != upload_root and not file_path.startswith(upload_root + os.sep):
        raise ValueError("非法的上传文件路径")
    return file_path if os.path.exists(file_path) else None


def _extract_document_pages(filename: str, contents: bytes) -> list[dict[str, Any]]:
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        return _extract_pdf_pages(contents)
    if ext == ".pptx":
        return _extract_pptx_pages(contents)
    if ext in {".txt", ".md"}:
        text = contents.decode("utf-8", errors="ignore").strip()
        return [{"page": 1, "text": text}] if text else []
    raise ValueError(f"暂不支持该文档格式: {ext}")


def _extract_pdf_pages(contents: bytes) -> list[dict[str, Any]]:
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise RuntimeError("缺少 PyMuPDF 依赖，请先安装 requirements.txt") from exc

    pages: list[dict[str, Any]] = []
    with fitz.open(stream=contents, filetype="pdf") as doc:
        for index, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                pages.append({"page": index, "text": _clean_text(text)})
    return pages


def _extract_pptx_pages(contents: bytes) -> list[dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("缺少 python-pptx 依赖，请先安装 requirements.txt") from exc

    prs = Presentation(io.BytesIO(contents))
    pages: list[dict[str, Any]] = []
    for index, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text).strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    rows.append(" | ".join(cell.text.strip() for cell in row.cells))
                if rows:
                    parts.append("\n".join(rows))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"备注：{notes}")
        text = _clean_text("\n\n".join(parts))
        if text:
            pages.append({"page": index, "text": text})
    return pages


def _clean_text(text: str) -> str:
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _format_pages_for_llm(pages: list[dict[str, Any]]) -> str:
    chunks = []
    for page in pages:
        chunks.append(f"=== 第 {page['page']} 页 ===\n{page['text']}")
    return "\n\n".join(chunks)


def _job_progress_updates(updates: dict[str, Any]) -> dict[str, Any]:
    """把内部进度字段映射到任务表字段。"""
    mapped: dict[str, Any] = {}
    if updates.get("page_count") is not None:
        mapped["page_count"] = str(updates["page_count"])
    if updates.get("text_chars") is not None:
        mapped["text_chars"] = str(updates["text_chars"])
    if updates.get("stage"):
        mapped["result"] = {
            "stage": updates.get("stage"),
            "chunk_index": updates.get("chunk_index"),
            "chunk_count": updates.get("chunk_count"),
        }
    if updates.get("extracted_text"):
        mapped.setdefault("result", {})
        mapped["result"]["extracted_text"] = updates["extracted_text"]
    return mapped


async def _extract_document_with_qwen(
    filename: str,
    pages: list[dict[str, Any]],
    progress_callback: Callable[[dict[str, Any]], Awaitable[None]] | None = None,
) -> dict[str, Any]:
    """根据文档大小选择直读或分块 ingest。"""
    raw_text = _format_pages_for_llm(pages)
    if len(pages) <= DOCUMENT_CHUNK_PAGES and len(raw_text) <= DOCUMENT_DIRECT_MAX_CHARS:
        if progress_callback:
            await progress_callback({"stage": "qwen_direct"})
        return await _extract_with_qwen(filename, raw_text)

    chunks = _split_page_chunks(pages, DOCUMENT_CHUNK_PAGES)
    chunk_results: list[dict[str, Any]] = []
    for index, chunk_pages in enumerate(chunks, start=1):
        if progress_callback:
            await progress_callback({
                "stage": "qwen_chunking",
                "chunk_index": index,
                "chunk_count": len(chunks),
            })
        chunk_text = _format_pages_for_llm(chunk_pages)
        chunk_results.append(await _extract_chunk_with_qwen(
            filename=filename,
            chunk_text=chunk_text,
            page_start=chunk_pages[0]["page"],
            page_end=chunk_pages[-1]["page"],
        ))

    if progress_callback:
        await progress_callback({
            "stage": "qwen_merging",
            "chunk_count": len(chunks),
        })
    return await _merge_chunk_results_with_qwen(filename, chunk_results)


def _split_page_chunks(pages: list[dict[str, Any]], chunk_size: int) -> list[list[dict[str, Any]]]:
    return [pages[i:i + chunk_size] for i in range(0, len(pages), chunk_size)]


async def _extract_chunk_with_qwen(
    filename: str,
    chunk_text: str,
    page_start: int,
    page_end: int,
) -> dict[str, Any]:
    system_prompt = (
        "你是文档分块信息抽取助手。"
        "只抽取当前页范围中明确出现的信息，不要补全或推测。"
        "请返回合法 JSON，供后续汇总成客户画像。"
    )
    user_prompt = (
        f"文件名：{filename}\n"
        f"页码范围：第 {page_start}-{page_end} 页\n\n"
        "请只返回合法 JSON，字段如下：\n"
        "{\n"
        '  "page_range": [1, 6],\n'
        '  "document_title_candidates": ["资料标题候选"],\n'
        '  "company_name_candidates": ["公司名称候选"],\n'
        '  "company_facts": ["公司/品牌/行业/业务/优势/渠道等事实，保留页码"],\n'
        '  "media_assets": [\n'
        '    {\n'
        '      "city": "城市",\n'
        '      "city_value": ["城市价值"],\n'
        '      "screen_location": "大屏位置",\n'
        '      "location_features": ["位置特点"],\n'
        '      "screen_features": ["大屏特点"],\n'
        '      "screen_specs": {"size": "", "resolution": "", "aspect_ratio": "", "orientation": "", "duration": "", "other_specs": []},\n'
        '      "daily_media_contacts": "日媒体接触人次/日均触达人次，如 30万人次/日",\n'
        '      "audience_or_traffic": ["人流/曝光/客群"],\n'
        '      "source_pages": [1]\n'
        '    }\n'
        '  ],\n'
        '  "requirements": ["需求/交付/周期/预算/限制等事实"],\n'
        '  "questions": ["需要确认的问题"],\n'
        '  "page_summaries": [{"page": 1, "title": "页标题", "summary": "页摘要", "facts": ["事实"]}]\n'
        "}\n\n"
        "当前分块正文：\n"
        f"{chunk_text}"
    )
    return await _post_qwen_json(system_prompt, user_prompt)


async def _merge_chunk_results_with_qwen(filename: str, chunk_results: list[dict[str, Any]]) -> dict[str, Any]:
    system_prompt = (
        "你是面向裸眼3D/数字艺术项目的客户资料汇总专家。"
        "请把多个分块抽取结果去重、合并成最终客户画像 JSON。"
        "必须忠实于分块结果；不确定的信息放入 questions 或 risks，不要编造。"
    )
    user_prompt = (
        f"文件名：{filename}\n\n"
        "请根据下面的分块抽取结果，汇总为最终 JSON。\n"
        f"{_result_schema_prompt()}\n\n"
        "分块抽取结果：\n"
        f"{json.dumps(chunk_results, ensure_ascii=False)}"
    )
    return await _post_qwen_json(system_prompt, user_prompt)


async def _extract_with_qwen(filename: str, raw_text: str) -> dict[str, Any]:
    if not settings.AI_API_KEY:
        raise RuntimeError("服务端未配置 AI_API_KEY")

    system_prompt = (
        "你是面向裸眼3D/数字艺术项目的资料 ingest 专家。"
        "请把管理员上传的 PDF/PPT 资料解析成可供后续 LLM 生成方案、派单和需求澄清使用的结构化 JSON。"
        "必须忠实于原文；不确定的信息放入 questions 或 risks，不要编造。"
    )
    user_prompt = (
        f"文件名：{filename}\n\n"
        f"{_result_schema_prompt()}\n\n"
        "文档正文：\n"
        f"{raw_text}"
    )
    return await _post_qwen_json(system_prompt, user_prompt)


def _result_schema_prompt() -> str:
    return (
        "请只返回合法 JSON，字段如下：\n"
        "{\n"
        '  "document_title": "资料标题",\n'
        '  "brief": "200字以内摘要",\n'
        '  "company_info": {\n'
        '    "company_name": "公司名称",\n'
        '    "brand_name": "品牌名/项目品牌",\n'
        '    "industry": "行业",\n'
        '    "business_scope": ["主营业务"],\n'
        '    "company_positioning": "公司定位",\n'
        '    "target_audience": ["目标客群"],\n'
        '    "products_or_services": ["核心产品/服务"],\n'
        '    "selling_points": ["品牌卖点/核心优势"],\n'
        '    "tone_and_style": ["品牌调性/视觉风格"],\n'
        '    "existing_assets": ["已有IP、素材、案例、视觉资产"],\n'
        '    "website_or_channels": ["官网、公众号、小红书、抖音等渠道"]\n'
        '  },\n'
        '  "media_assets": [\n'
        '    {\n'
        '      "city": "大屏所在城市",\n'
        '      "city_value": ["城市商业价值、人流价值、消费价值、地标价值、传播价值"],\n'
        '      "screen_location": "大屏具体位置/商圈/楼体/广场",\n'
        '      "location_features": ["该位置的特点，如核心商圈、交通枢纽、游客聚集、夜经济、人群画像"],\n'
        '      "screen_features": ["大屏特点，如裸眼3D、L型屏、曲面屏、超高清、户外LED、沉浸式"],\n'
        '      "screen_specs": {\n'
        '        "size": "尺寸/面积，如 800㎡",\n'
        '        "resolution": "分辨率",\n'
        '        "aspect_ratio": "画幅比例",\n'
        '        "orientation": "横屏/竖屏/转角/L型/曲面",\n'
        '        "duration": "常规播放时长/素材时长",\n'
        '        "other_specs": ["亮度、点间距、播放频次等其他规格"]\n'
        '      },\n'
        '      "daily_media_contacts": "日媒体接触人次/日均触达人次，如 30万人次/日",\n'
        '      "audience_or_traffic": ["人流量、客群、曝光量、消费层级等"],\n'
        '      "source_pages": [1]\n'
        '    }\n'
        '  ],\n'
        '  "project_requirements": ["明确需求、目标、受众、场景"],\n'
        '  "creative_direction": ["视觉风格、创意方向、参考案例"],\n'
        '  "deliverables": ["交付物、格式、数量、规格"],\n'
        '  "timeline_budget": ["周期、节点、预算、报价相关信息"],\n'
        '  "risks": ["限制、矛盾、不完整或需要注意的信息"],\n'
        '  "questions": ["下一步需要向客户确认的问题"],\n'
        '  "page_summaries": [{"page": 1, "title": "页标题", "summary": "页摘要", "facts": ["带页码来源的事实"]}]\n'
        "}"
    )


async def _post_qwen_json(system_prompt: str, user_prompt: str) -> dict[str, Any]:
    if not settings.AI_API_KEY:
        raise RuntimeError("服务端未配置 AI_API_KEY")
    base_url = (settings.DOCUMENT_INGEST_BASE_URL or settings.AI_BASE_URL).rstrip("/")
    model = settings.DOCUMENT_INGEST_MODEL_NAME or settings.AI_MODEL_NAME
    timeout_seconds = float(settings.DOCUMENT_INGEST_HTTP_TIMEOUT or 240)
    try:
        async with httpx.AsyncClient(timeout=timeout_seconds) as client:
            response = await client.post(
                f"{base_url}/chat/completions",
                headers={
                    "Authorization": f"Bearer {settings.AI_API_KEY}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": model,
                    "messages": [
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt},
                    ],
                    "response_format": {"type": "json_object"},
                },
            )
            response.raise_for_status()
            content = response.json()["choices"][0]["message"]["content"]
            return _parse_json_object(content)
    except httpx.TimeoutException as exc:
        raise RuntimeError(f"Qwen 文档解析超时（{timeout_seconds:.0f}s），文档可能较长或模型响应较慢") from exc
    except httpx.HTTPStatusError as exc:
        detail = exc.response.text[:500] if exc.response is not None else ""
        raise RuntimeError(f"Qwen 文档解析接口返回错误: {exc.response.status_code} {detail}") from exc


def _parse_json_object(content: str) -> dict[str, Any]:
    content = content.strip()
    block_match = re.search(r"```(?:json)?\s*([\s\S]*?)```", content)
    if block_match:
        content = block_match.group(1).strip()
    if not content.startswith("{"):
        object_match = re.search(r"\{[\s\S]*\}", content)
        if object_match:
            content = object_match.group(0)
    data = json.loads(content)
    if not isinstance(data, dict):
        raise ValueError("LLM 返回的不是 JSON object")
    return data


def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def _filename_stem(filename: str) -> str:
    stem = os.path.splitext(os.path.basename(filename or "document"))[0].strip()
    return re.sub(r"[\\/:*?\"<>|]+", "_", stem) or "document"
